"""Build final manuscript .docx (JA + EN) with inline figures/tables
and a single editable .pptx (EN figures) for per-slide editing.

Formatted for Economic Analysis and Policy / Elsevier journals:
  - Title page with Declarations uploaded separately for double-anonymized review
  - 1.5× line spacing, 10pt Times New Roman body text
  - Decimal heading system (1, 1.1, 1.1.1) with max three levels
  - Author-date (Harvard) references: sentence-case surnames, no ampersand,
    DOIs as full links
  - Centred page numbers in footer
  - Figures inline in manuscript + separate editable .pptx and PNG files

Usage:  python build_docx_pptx.py
Outputs into ../manuscript/ :
  - manuscript_en.docx / .pdf
  - manuscript_ja.docx / .pdf
  - supporting_information_en.docx / .pdf
  - figures_en.pptx
  - table*.docx
  - title_page_en.docx / .pdf
  - highlights_en.docx / .pdf
  - cover_letter_en.docx / .pdf
"""
from __future__ import annotations

import os
import re
import json
import glob
import shutil
import functools
from dataclasses import dataclass

import pandas as pd
from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from latex2word import LatexToWordElement
import lxml.etree as ET

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
MS = os.path.join(ROOT, "manuscript")
FIG = os.path.join(ROOT, "figures")
TAB = os.path.join(ROOT, "tables")
DATA = os.path.join(ROOT, "data")
os.makedirs(MS, exist_ok=True)

# Maps from content key to (prefix, display_number), populated while building the
# manuscript so that the PPTX and separate figure files match the inline numbering.
FIGURE_DISPLAY_MAP = {}
TABLE_DISPLAY_MAP = {}

# ---------------------------------------------------------------------------
# Figure/table metadata

FIG_LIST = [
    ("fig1", "Fig. 1", "図1",
     "In-sample 1-year GDP growth RMSE across M0–M4 for 39 countries (lower is better).",
     "39カ国における単年GDP成長率RMSEの標本内比較（M0–M4、小さいほど良い）。",
     "fig1_m_ranking_{lang}.png"),
    ("fig2", "Fig. 2", "図2",
     "Out-of-sample MAPE on 2015–19 held-out window; M2 achieves 13% relative improvement vs M0.",
     "2015–19年のホールドアウト窓における標本外MAPE。M2はM0比で相対13%改善。",
     "fig2_oos_{lang}.png"),
    ("fig3", "Fig. 3", "図3",
     "PIM tangible-plus-intangible stock versus CWON PCA (within-country demeaned log).",
     "PIM有形・無形資本ストックとCWON PCAの軌跡比較（国内平均除去対数）。",
     "fig3_trajectories_{lang}.png"),
    ("fig4", "Fig. 4", "図4",
     "Price-sensitivity of the PIM/CWON log-ratio. Left: selected countries. Right: cross-country zero-crossing γ_price; all roots lie outside the explored [-0.04, +0.04] range.",
     "PIM/CWON対数比の価格感度。左：代表国。右：国別零点γ_price；全ての零点は探索範囲[-0.04, +0.04]外にある。",
     "fig4_gamma_price_{lang}.png"),
    ("fig5", "Fig. 5", "図5",
     "Conceptual diagram of the population-capital tempo correspondence.",
     "人口・資本テンポ対応関係の概念図。",
     "fig5_concept_{lang}.png"),
    ("fig6", "Fig. 6", "図6",
     "Relational PIM diagnostics: ρ̂₂ across 39 countries under M0 vs M4.",
     "関係型PIM診断: M0とM4における39カ国のρ̂₂。",
     "fig6_rpim_{lang}.png"),
    ("fig7", "Fig. 7", "図7",
     "Depreciation–lag sensitivity: estimated μ̂ under ±20% depreciation perturbations.",
     "減価償却率–ラグ感度: ±20%変動下の推定μ̂。",
     "fig7_delta_sensitivity_{lang}.png"),
    ("fig8", "Fig. 8", "図8",
     "Conditional OOS evaluation: interior-solution vs boundary countries.",
     "条件付き標本外評価: 内点解国 vs 境界解国。",
     "fig8_conditional_oos_{lang}.png"),
    ("fig9", "Fig. 9", "図9",
     "Cross-sectional regression of ρ̂₂ on R&D intensity.",
     "ρ̂₂のR&D強度に対するクロスセクション回帰。",
     "fig9_rho2_regression_{lang}.png"),
    ("fig10", "Fig. 10", "図10",
     "Capital-stock divergence: observed versus baseline capital stock over time for six representative countries.",
     "資本ストック乖離: 6カ国における観測値とベースライン資本ストックの推移。",
     "fig10_k_divergence_{lang}.png"),
    ("fig11", "Fig. 11", "図11",
     "Measurement consequence: K-level change vs log-level TFP shift (country means, 2010–2019).",
     "計測帰結: K水準変化と対数水準TFPシフト（国別平均、2010–2019）。",
     "fig11_tfp_consequence_{lang}.png"),
    ("fig12", "Fig. 12", "図12",
     "Implied labour-share correction from tempo-adjusted capital (country means, 2010–2019).",
     "テンポ調整資本からの労働分配率補正（国別平均、2010–2019）。",
     "fig12_labor_share_{lang}.png"),
    ("fig13", "Fig. 13", "図13",
     "Solow-residual decomposition: M0 vs tempo-adjusted (M2) vs joint (M4) for six representative countries.",
     "ソロー残差の分解: M0 vs テンポ調整(M2) vs 統合(M4)、代表6カ国。",
     "fig13_solow_decomp_{lang}.png"),
    ("fig14", "Fig. 14", "図14",
     "National wealth: CWON official vs intangible-adjusted produced capital (2019).",
     "国富: CWON公式値 vs 無形資本調整後の生産資本（2019年）。",
     "fig14_counterfactual_wealth_{lang}.png"),
    ("fig15", "Fig. 15", "図15",
     "2020–2040 GDP-level scenarios (2019=100) for selected countries, including M4 baseline, level-only TFP (2019 held fixed), AI surge, and fiscal-stimulus variants.",
     "2020–2040年のGDP水準シナリオ（2019年=100）：代表国別に、M4ベースライン、TFPレベル固定、AI急増、財政刺激のバリエーションを含む。",
     "fig15_future_scenarios_{lang}.png"),
    ("fig16", "Fig. 16", "図16",
     "Country clusters by joint-identified (μ̂, β̂) and asset composition; silhouette and bootstrap-stability metrics support the k=3 partition (Table 10b).",
     "(μ̂, β̂) と資産構成による国クラスター；シルエット・ブートストラップ安定性指標がk=3の分割を支持する（Table 10b）。",
     "fig16_country_clusters_{lang}.png"),
    ("fig17", "Fig. 17", "図17",
     "Time-varying depreciation robustness: median estimated lag.",
     "時変減価償却率の頑健性: 推定ラグの中央値。",
     "fig17_delta_timevarying_{lang}.png"),
    ("fig18", "Fig. 18", "図18",
     "Monte Carlo identification sharpness for four calibrated economies (United States, Republic of Korea, France, Colombia). Error bars show approximate 95 % confidence intervals for the mean RMSE across parameter cells.",
     "4カ国（米国、韓国、フランス、コロンビア）のモンテカルロ同定鋭度。誤差線は各パラメータセルのモンテカルロ標準誤差から導出した平均RMSEのおおよその95%信頼区間を示す。",
     "fig18_monte_carlo_{lang}.png"),
]


# ---------------------------------------------------------------------------
# docx helpers

def set_font(run, name="Times New Roman", size=10, bold=False, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.bold = bold
    run.italic = italic
    # Apply East Asian font for proper JA rendering
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:ascii"), name)
    rFonts.set(qn("w:hAnsi"), name)
    rFonts.set(qn("w:eastAsia"), "MS Mincho")


def format_jpa_heading(text: str, level: int,
                        section_num: int = 0,
                        subsection_num: int = 0) -> str:
    """Return Journal of Productivity Analysis style heading text.

    - Level 1 (# in md): paper title — title case, centered.
    - Level 2 (## in md): section heading — decimal number + sentence case.
    - Level 3 (### in md): subsection heading — decimal prefix + sentence case.
    """
    # Strip leading decimal marker if present ("1 Introduction" → "Introduction")
    clean = re.sub(r'^\d+(\.\d+)?\s+', '', text)
    if level == 1:
        return clean
    if level == 2:
        if section_num == 0:
            # Special headings like "References" / "Tables" remain unnumbered
            return clean
        return f"{section_num} {clean}"
    # level 3
    if section_num and subsection_num:
        return f"{section_num}.{subsection_num} {clean}"
    return clean


def add_heading(doc, text, level, lang, section_num=0, subsection_num=0):
    """Add heading with Economica formatting.

    Level 1 (# in md): paper title — centered, 16pt bold, title case.
    Level 2 (## in md): section heading — centered, 14pt bold, ALL CAPS, Roman numeral.
    Level 3 (### in md): subsection — flush left, 12pt bold, Roman prefix.
    """
    h = doc.add_paragraph()
    h.paragraph_format.space_before = Pt(18)
    h.paragraph_format.space_after = Pt(6)
    h.paragraph_format.line_spacing = 1.5
    h.paragraph_format.first_line_indent = Pt(0)

    display_text = format_jpa_heading(text, level, section_num, subsection_num)

    # Title: centered; section/subsection headings: left aligned
    if level == 1:
        h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    else:
        h.alignment = WD_ALIGN_PARAGRAPH.LEFT

    run = h.add_run(display_text)
    set_font(run, size={1: 16, 2: 12, 3: 12}[level], bold=True)
    return h


# Cross-reference patterns: Arabic section/subsection numbers → Roman
SECTION_REF_RE = re.compile(
    r'\b(Sect(?:ion|s)?\.?|§)\s*(\d+)(?:\.(\d+))?',
    re.IGNORECASE,
)

def preserve_en_crossrefs(text: str) -> str:
    """Leave decimal section references unchanged (JPA style)."""
    return text


def format_elsevier_reference(line: str) -> str:
    """Format a reference line for Elsevier Economic Modelling (Harvard style).

    The reference source is already in author-date form with journal/book titles
    marked in markdown italics.  Keep the markdown markup so the Word renderer
    can apply italics, and only do light clean-up (strip stray bullets, collapse
    whitespace, and remove spurious emphasis on the whole line).
    """
    line = line.strip()
    if not line:
        return line
    # Drop accidental leading list markers
    if line.startswith(("- ", "* ")) and not line.startswith("**"):
        line = line[2:].strip()
    # Collapse multiple spaces, but keep the markdown emphasis and en-dashes
    line = re.sub(r"[ \t]+", " ", line)
    return line


def int_to_roman(n: int) -> str:
    values = [
        (1000, "M"), (900, "CM"), (500, "D"), (400, "CD"),
        (100, "C"), (90, "XC"), (50, "L"), (40, "XL"),
        (10, "X"), (9, "IX"), (5, "V"), (4, "IV"), (1, "I"),
    ]
    out = []
    for v, s in values:
        while n >= v:
            out.append(s)
            n -= v
    return "".join(out)


GREEK_LETTERS = set("αβγδεζηθικλμνξοπρσςτυφχψω"
                     "ΑΒΓΔΕΖΗΘΙΚΛΜΝΞΟΠΡΣΤΥΦΧΨΩ")

# Unicode math → LaTeX command map (used for Word OMML equations)
GREEK_TO_LATEX = {
    'α': '\\alpha', 'β': '\\beta', 'γ': '\\gamma', 'δ': '\\delta',
    'ε': '\\varepsilon', 'θ': '\\theta', 'λ': '\\lambda', 'μ': '\\mu',
    'ρ': '\\rho', 'σ': '\\sigma', 'Σ': '\\sum', 'π': '\\pi',
    'Δ': '\\Delta', 'ς': '\\varsigma', 'ζ': '\\zeta', 'η': '\\eta',
    'ι': '\\iota', 'κ': '\\kappa', 'ν': '\\nu', 'ξ': '\\xi',
    'ο': '\\omicron', 'τ': '\\tau', 'υ': '\\upsilon', 'φ': '\\phi',
    'χ': '\\chi', 'ψ': '\\psi', 'ω': '\\omega',
    'Α': '\\mathrm{A}', 'Β': '\\mathrm{B}', 'Γ': '\\Gamma', 'Δ': '\\Delta',
    'Ε': '\\mathrm{E}', 'Ζ': '\\mathrm{Z}', 'Η': '\\mathrm{H}', 'Θ': '\\Theta',
    'Ι': '\\mathrm{I}', 'Κ': '\\mathrm{K}', 'Λ': '\\Lambda', 'Μ': '\\mathrm{M}',
    'Ν': '\\mathrm{N}', 'Ξ': '\\Xi', 'Ο': '\\mathrm{O}', 'Π': '\\Pi',
    'Ρ': '\\mathrm{P}', 'Σ': '\\Sigma', 'Τ': '\\mathrm{T}', 'Υ': '\\Upsilon',
    'Φ': '\\Phi', 'Χ': '\\mathrm{X}', 'Ψ': '\\Psi', 'Ω': '\\Omega',
}
SUBSCRIPT_UNICODE = {
    '₀': '0', '₁': '1', '₂': '2', '₃': '3', '₄': '4', '₅': '5',
    '₆': '6', '₇': '7', '₈': '8', '₉': '9', 'ₐ': 'a', 'ₑ': 'e',
    'ₕ': 'h', 'ᵢ': 'i', 'ⱼ': 'j', 'ₖ': 'k', 'ₗ': 'l', 'ₘ': 'm',
    'ₙ': 'n', 'ₒ': 'o', 'ₚ': 'p', 'ᵣ': 'r', 'ₛ': 's', 'ₜ': 't',
    'ᵤ': 'u', 'ᵥ': 'v', 'ᵥ': 'v', 'ₓ': 'x', 'ᵦ': '\\beta',
    'ᵧ': '\\gamma', 'ᵨ': '\\rho', 'ᵩ': '\\phi', 'ᵪ': '\\chi',
}
SUPERSCRIPT_UNICODE = {
    '²': '2', '³': '3', '⁴': '4', '⁵': '5', '⁶': '6', '⁷': '7',
    '⁸': '8', '⁹': '9', '⁰': '0', '¹': '1',
    'ᵃ': 'a', 'ᵇ': 'b', 'ᶜ': 'c', 'ᵈ': 'd', 'ᵉ': 'e', 'ᶠ': 'f',
    'ᵍ': 'g', 'ʰ': 'h', 'ⁱ': 'i', 'ʲ': 'j', 'ᵏ': 'k', 'ˡ': 'l',
    'ᵐ': 'm', 'ⁿ': 'n', 'ᵒ': 'o', 'ᵖ': 'p', 'ʳ': 'r', 'ˢ': 's',
    'ᵗ': 't', 'ᵘ': 'u', 'ᵛ': 'v', 'ʷ': 'w', 'ˣ': 'x', 'ʸ': 'y', 'ᶻ': 'z',
}

# Combining diacritics treated as LaTeX accents
COMBINING_ACCENTS = {
    '\u0302': 'hat',   # circumflex (e.g. μ̂)
    '\u030c': 'check', # caron
    '\u0306': 'breve', # breve
}

# Tokens that should be rendered in italic when they appear as standalone
# variable-like identifiers in math context.
MATH_VARS = re.compile(
    r"""(?<![A-Za-z_])      # not preceded by letter/underscore
    (                        # group 1: the variable token
      [A-Z]                  # single capital letter ...
      (?:_[A-Za-z0-9]+)?    # ... optionally followed by _subscript
      |                      # OR
      [αβγδεζηθικλμνξοπρσςτυφχψωΑΒΓΔΕΖΗΘΙΚΛΜΝΞΟΠΡΣΤΥΦΧΨΩ] # Greek letter
      (?:_[A-Za-z0-9]+)?    # optional subscript
    )
    (?![A-Za-z])             # not followed by letter
    """,
    re.VERBOSE,
)

# Pattern to detect subscript notations: _{...} or _X (single char)
SUB_RE = re.compile(r'_\{([^}]+)\}|_([A-Za-z0-9])')
# Pattern to detect superscript notations: ^{...} or ^X (single char) or trailing *
SUP_RE = re.compile(r'\^\{([^}]+)\}|\^([A-Za-z0-9])')

# Displayed equation: 4-space indent, optional label like (M0), (1), (2)
EQUATION_RE = re.compile(r'^    (.+?)\s{2,}\(([A-Za-z0-9]+)\)\s*$')
EQUATION_NOLABEL_RE = re.compile(r'^    (.+?)\s*$')


def _is_greek(ch):
    return ch in GREEK_LETTERS


def _set_table_borders_none(table):
    """Remove all table cell borders."""
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        return
    tblBorders = tblPr.first_child_found_in('w:tblBorders')
    if tblBorders is None:
        tblBorders = OxmlElement('w:tblBorders')
        tblPr.append(tblBorders)
    for edge in ('top', 'left', 'bottom', 'right', 'insideH', 'insideV'):
        edge_el = tblBorders.find(qn(f'w:{edge}'))
        if edge_el is None:
            edge_el = OxmlElement(f'w:{edge}')
            tblBorders.append(edge_el)
        edge_el.set(qn('w:val'), 'none')
        edge_el.set(qn('w:sz'), '0')
        edge_el.set(qn('w:space'), '0')
        edge_el.set(qn('w:color'), 'auto')


def unicode_math_to_latex(s):
    """Convert the Unicode math strings used in the source markdown to LaTeX.

    Handles Greek letters, Unicode sub/superscripts, operators, and
    multi-character subscripts so that latex2mathml/mathml2omml can turn the
    result into a native Word OMML equation.
    """
    if not s:
        return s

    out = []
    i = 0
    while i < len(s):
        ch = s[i]
        # Combining circumflex/carons on Greek or Latin letters -> \hat{}, etc.
        if i + 1 < len(s) and s[i + 1] in COMBINING_ACCENTS and (ch.isalpha() or _is_greek(ch)):
            accent = COMBINING_ACCENTS[s[i + 1]]
            base = GREEK_TO_LATEX[ch] if _is_greek(ch) else ch
            out.append(' \\{}{{{}}} '.format(accent, base))
            i += 2
            continue
        if ch in SUBSCRIPT_UNICODE:
            out.append('_' + SUBSCRIPT_UNICODE[ch])
        elif ch in SUPERSCRIPT_UNICODE:
            out.append('^' + SUPERSCRIPT_UNICODE[ch])
        elif ch in GREEK_TO_LATEX:
            if ch == 'Σ' and i + 1 < len(s) and s[i + 1] in SUBSCRIPT_UNICODE:
                sub = SUBSCRIPT_UNICODE[s[i + 1]]
                out.append('\\sum_{' + sub + '}')
                i += 2
                continue
            # Use \\mathrm for capital Greek that are identifiers; otherwise
            # the LaTeX command name with surrounding spaces.
            out.append(' ' + GREEK_TO_LATEX[ch] + ' ')
        elif ch == '−':
            out.append('-')
        elif ch == '·':
            out.append(' \\cdot ')
        elif ch == '×':
            out.append(' \\times ')
        elif ch == '∈':
            out.append(' \\in ')
        elif ch == '≈':
            out.append(' \\approx ')
        elif ch == '≤':
            out.append(' \\leq ')
        elif ch == '≥':
            out.append(' \\geq ')
        elif ch == '≠':
            out.append(' \\neq ')
        else:
            out.append(ch)
        i += 1

    s2 = ''.join(out)
    # Convert Sigma with explicit subscript to \sum_{...}
    s2 = re.sub(r'\\sum\s+_([A-Za-z0-9])', r'\\sum_{\1}', s2)
    s2 = re.sub(r'\\sum\s+_\{([^}]+)\}', r'\\sum_{\1}', s2)
    # Wrap multi-letter / digit subscripts: K_tang -> K_{tang}, Y_t -> Y_{t}
    # Avoid touching already-braced subscripts like _{t-1-s}.
    s2 = re.sub(r'(?<!_)_([A-Za-z][A-Za-z0-9]*)', r'_{\1}', s2)
    # Render common function names as upright operators
    for fn in ('log', 'exp', 'sin', 'cos', 'min', 'max'):
        s2 = re.sub(rf'\b{fn}\b', rf'\\{fn}', s2)
    # Treat the fertility-rate acronym TFR as a text label in equations
    s2 = re.sub(r'\bTFR\b', r'\\text{TFR}', s2)
    # Tighten sub/superscript braces next to closing braces
    s2 = re.sub(r'(?<=\})\s+_', '_', s2)
    s2 = re.sub(r'(?<=\})\s+\^', '^', s2)
    # Collapse multiple spaces (LaTeX tokeniser is tolerant but keeps output tidy)
    s2 = re.sub(r'  +', ' ', s2).strip()
    return s2


def add_omml_math(paragraph, latex, size=10):
    """Append a native Word OMML equation from a LaTeX string to a paragraph."""
    try:
        elem = LatexToWordElement(latex).element()
        _fix_omml_accents(elem)
        paragraph._element.append(elem)
    except Exception:
        # Fallback: render the raw text as styled runs so the build never breaks.
        _add_run(paragraph, latex, size=size, italic=True)


_MNS = '{http://schemas.openxmlformats.org/officeDocument/2006/math}'
_OMML_ACCENTS = set('^~¯˙˚ˇˆ˜´`\u0302\u0304\u2192\u20d7')

def _fix_omml_accents(elem):
    """Convert OMML <m:limUpp> accent fakes to proper <m:acc> elements.

    latex2mathml emits \hat / \bar as <mover> with a plain accent character.
    mathml2omml turns those into <m:limUpp> blocks, which Word renders as
    superscripts rather than accents. Replace them with <m:acc> when the
    upper "limit" is a single known accent character.
    """
    if elem is None:
        return
    for limupp in list(elem.iter(_MNS + 'limUpp')):
        base = None
        lim = None
        for child in limupp:
            if child.tag == _MNS + 'e' and base is None:
                base = child
            elif child.tag == _MNS + 'lim':
                lim = child
        if base is None or lim is None:
            continue
        # Extract the accent text from all descendant m:t elements
        accent_parts = []
        for t in lim.iter(_MNS + 't'):
            if t.text:
                accent_parts.append(t.text)
        accent = ''.join(accent_parts).strip()
        if accent not in _OMML_ACCENTS or len(accent) != 1:
            continue
        # Unwrap a redundant m:box around the base if present
        source = base
        if len(base) == 1 and base[0].tag == _MNS + 'box':
            for box_child in base[0]:
                if box_child.tag == _MNS + 'e':
                    source = box_child
                    break
        # Build <m:acc><m:accPr><m:chr m:val="..."/></m:accPr><m:e>...</m:e></m:acc>
        new_acc = ET.Element(_MNS + 'acc')
        acc_pr = ET.SubElement(new_acc, _MNS + 'accPr')
        chr_el = ET.SubElement(acc_pr, _MNS + 'chr')
        chr_el.set(_MNS + 'val', accent)
        new_e = ET.SubElement(new_acc, _MNS + 'e')
        for child in list(source):
            new_e.append(child)
        parent = limupp.getparent()
        if parent is None:
            continue
        idx = list(parent).index(limupp)
        parent[idx] = new_acc


def _add_run(paragraph, text, size=10, bold=False, italic=False,
             superscript=False, subscript=False, font_name="Times New Roman"):
    """Add a single run with specified formatting."""
    run = paragraph.add_run(text)
    set_font(run, name=font_name, size=size, bold=bold, italic=italic)
    if superscript:
        run.font.superscript = True
    if subscript:
        run.font.subscript = True
    return run


def _split_italic(text):
    """Split text into (is_italic, content) segments respecting *...* markers.
    Asterisks inside {} braces or preceded by a backslash are not treated as
    italic markers so that LaTeX-style *X^{*}* can be used inline."""
    segments = []
    i = 0
    n = len(text)
    cur = []
    in_italic = False
    depth = 0
    while i < n:
        ch = text[i]
        if ch == '\\' and i + 1 < n and text[i + 1] in '*_\\`':
            cur.append(text[i + 1])
            i += 2
            continue
        if ch == '*':
            if depth == 0:
                if cur:
                    segments.append((in_italic, ''.join(cur)))
                in_italic = not in_italic
                cur = []
                i += 1
                continue
        if ch == '{':
            depth += 1
        elif ch == '}':
            depth = max(0, depth - 1)
        cur.append(ch)
        i += 1
    if cur:
        segments.append((in_italic, ''.join(cur)))
    return segments


_MATH_SYMBOLS_RE = re.compile(
    r'[\u0370-\u03ff\u1f00-\u1fff]|_|=|/|−|-|\*|·|×|∈|≠|≈|≤|≥|\(|\)|\{|\}|\^|[\u2070-\u207f\u2080-\u209c]')


def _is_math_span(text):
    """Return True if *text* should be rendered as a native Word OMML equation."""
    stripped = text.strip()
    if not stripped:
        return False

    if _MATH_SYMBOLS_RE.search(stripped):
        return True
    # Uppercase/digit identifiers such as M0, M_obs, K_tang.
    if re.fullmatch(r'[A-Z][A-Za-z0-9]*(?:_[A-Za-z0-9]+)?', stripped) and ('_' in stripped or any(c.isdigit() for c in stripped)):
        return True
    return False


def add_math_runs(paragraph, text, size=10, base_italic=False):
    """Parse *text* and emit Word runs or OMML with proper math formatting.

    *...* italic spans that contain math symbols are converted to native OMML
    equations; all other spans are emitted as styled runs.
    """
    if not text:
        return

    for is_italic, content in _split_italic(text):
        if not content:
            continue
        if is_italic and _is_math_span(content):
            latex = unicode_math_to_latex(content)
            add_omml_math(paragraph, latex, size)
        else:
            _emit_math_segment(paragraph, content, size,
                               force_italic=is_italic or base_italic)


def _extract_base_math_token(s, i):
    """Return a base math token starting at index i, or None."""
    n = len(s)
    if i >= n:
        return None
    ch = s[i]
    end = i
    if ch in GREEK_LETTERS:
        end = i + 1
        if end < n and s[end] in COMBINING_ACCENTS:
            end += 1
    elif ch.isupper():
        # Uppercase Latin identifiers: require at least one underscore so
        # model labels such as M0-M4 are not converted to equations.
        j = i
        while j < n and (s[j].isalnum() or s[j] == '_'):
            j += 1
        token = s[i:j]
        if '_' not in token:
            return None
        end = j
    else:
        return None

    # Consume trailing subscripts/superscripts and parenthesised arguments
    while end < n:
        if s[end] in SUBSCRIPT_UNICODE or s[end] in SUPERSCRIPT_UNICODE:
            end += 1
            continue
        if s[end] == '_':
            if end + 1 < n and s[end + 1] == '{':
                k = s.find('}', end + 2)
                if k != -1:
                    end = k + 1
                    continue
            k = end + 1
            while k < n and (s[k].isalnum() or s[k] == '_'):
                k += 1
            if k == end + 1:
                break
            end = k
            continue
        if s[end] == '^':
            if end + 1 < n and s[end + 1] == '{':
                k = s.find('}', end + 2)
                if k != -1:
                    end = k + 1
                    continue
            k = end + 1
            while k < n and (s[k].isalnum() or s[k] in '+-_'):
                k += 1
            if k == end + 1:
                break
            end = k
            continue
        if s[end] == '(':
            depth = 1
            k = end + 1
            while k < n and depth > 0:
                if s[k] == '(':
                    depth += 1
                elif s[k] == ')':
                    depth -= 1
                k += 1
            if depth == 0:
                end = k
                continue
            break
        break
    return s[i:end]


def _is_math_expression_char(ch):
    """Return True if ch can appear inside an inline math expression."""
    if ch.isspace():
        return True
    if ch.isdigit() or ch in '.%':
        return True
    if ch in '=+-−-*/×·^_%≠≈≤≥':
        return True
    if ch in GREEK_LETTERS:
        return True
    if ch in SUBSCRIPT_UNICODE or ch in SUPERSCRIPT_UNICODE:
        return True
    return False


def _extract_extended_math_token(s, i):
    """Return the longest inline math token starting at i, or None."""
    base = _extract_base_math_token(s, i)
    if not base:
        return None
    n = len(s)
    end = i + len(base)
    while end < n:
        if s[end].isspace():
            end += 1
            continue
        next_token = _extract_base_math_token(s, end)
        if next_token:
            end += len(next_token)
            continue
        if _is_math_expression_char(s[end]):
            end += 1
            continue
        break
    return s[i:end]


def _emit_math_segment(paragraph, inner, size, force_italic=False,
                       force_bold=False):
    """Emit a text segment with math-aware formatting.

    Greek letters and variable-like identifiers are converted to native Word
    OMML equations; everything else is emitted as styled runs.
    """
    if force_bold:
        _add_run(paragraph, inner, size=size, bold=True, italic=force_italic)
        return

    pos = 0
    n = len(inner)
    text_buf = []
    while pos < n:
        token = _extract_extended_math_token(inner, pos)
        if token:
            if text_buf:
                _add_run(paragraph, ''.join(text_buf), size=size,
                         italic=force_italic)
                text_buf = []
            math_token = token.rstrip()
            trailing = token[len(math_token):]
            latex = unicode_math_to_latex(math_token)
            add_omml_math(paragraph, latex, size)
            if trailing:
                text_buf.append(trailing)
            pos += len(token)
        else:
            text_buf.append(inner[pos])
            pos += 1
    if text_buf:
        _add_run(paragraph, ''.join(text_buf), size=size,
                 italic=force_italic)


def add_para(doc, text, lang, italic=False):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.first_line_indent = Pt(24)
    add_math_runs(p, text, size=10, base_italic=italic)
    return p


def add_rich_para(doc, text, lang, bullet=False):
    """Add a paragraph with bold **...** spans and math-aware formatting."""
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.5
    if bullet:
        p.paragraph_format.left_indent = Pt(36)
        p.paragraph_format.first_line_indent = Pt(-18)
    else:
        p.paragraph_format.first_line_indent = Pt(24)

    # Split on **bold** markers first
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**') and len(part) > 4:
            inner = part[2:-2]
            # Bold span — emit with bold + math formatting
            _emit_bold_math(p, inner, size=10)
        else:
            # Normal text — emit with math formatting (handles *italic*, Greek, sub/sup)
            add_math_runs(p, part, size=10, base_italic=False)
    return p


def _emit_bold_math(paragraph, text, size=10):
    """Emit bold text with math-aware formatting (Greek italic, sub/sup)."""
    for is_italic, content in _split_italic(text):
        if not content:
            continue
        _emit_math_segment(paragraph, content, size,
                           force_italic=is_italic, force_bold=True)


def add_equation_block(doc, equation_text, label=None):
    """Add a displayed equation as a native Word OMML object.

    The equation is rendered on the left of a borderless two-cell table and
    the label (e.g. (M0)) is right-aligned in the second cell. This keeps the
    label flush right without relying on tab stops inside an OMML paragraph.
    """
    latex = unicode_math_to_latex(equation_text.strip())

    # Some displayed equations use slash as a fraction; convert the first
    # left/right-variable pair into \frac to produce a proper OMML fraction.
    latex = re.sub(
        r'([A-Za-z][A-Za-z0-9_{}\\]*)\s*/\s*([A-Za-z][A-Za-z0-9_{}\\]*)',
        r'\\frac{\1}{\2}',
        latex,
        count=1,
    )

    table = doc.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    table.allow_autofit = False
    _set_table_borders_none(table)

    # Set column widths: ~16 cm total page width minus margins
    total_width = Inches(6.3)
    label_width = Inches(0.6)
    eq_width = Inches(total_width.inches - label_width.inches)
    table.columns[0].width = eq_width
    table.columns[1].width = label_width

    eq_cell = table.cell(0, 0)
    label_cell = table.cell(0, 1)

    # Equation cell: centred, no first-line indent
    eq_para = eq_cell.paragraphs[0]
    eq_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    eq_para.paragraph_format.space_before = Pt(12)
    eq_para.paragraph_format.space_after = Pt(12)
    eq_para.paragraph_format.line_spacing = 1.5
    eq_para.paragraph_format.first_line_indent = Pt(0)
    add_omml_math(eq_para, latex, size=10)

    # Label cell: right-aligned, vertically centered
    label_para = label_cell.paragraphs[0]
    label_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    label_para.paragraph_format.space_before = Pt(12)
    label_para.paragraph_format.space_after = Pt(12)
    label_para.paragraph_format.line_spacing = 1.5
    label_para.paragraph_format.first_line_indent = Pt(0)
    if label:
        label_run = label_para.add_run(f"({label})")
        set_font(label_run, size=10, italic=False)


def add_page_number_footer(section):
    """Add a centred page-number field to the section footer."""
    footer = section.footer
    p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(0)
    p.paragraph_format.space_before = Pt(0)
    p.text = ""
    fld = OxmlElement('w:fldSimple')
    fld.set(qn('w:instr'), 'PAGE \\* MERGEFORMAT')
    r = OxmlElement('w:r')
    t = OxmlElement('w:t')
    t.text = '1'
    r.append(t)
    fld.append(r)
    p._p.append(fld)


def add_figure(doc, png_path, caption_prefix, caption_text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run()
    run.add_picture(png_path, width=Inches(6.2))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap.paragraph_format.space_after = Pt(12)
    r1 = cap.add_run(f"{caption_prefix}. ")
    set_font(r1, size=10, bold=True)
    r2 = cap.add_run(caption_text)
    set_font(r2, size=10)


def add_dataframe_as_table(doc, df: pd.DataFrame, col_widths=None, font_size=10):
    df = df.fillna('')
    tbl = doc.add_table(rows=1 + len(df), cols=len(df.columns))
    # Remove any table style so direct tblBorders are not overridden
    tbl_element = tbl._tbl
    tblPr = tbl_element.tblPr
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr')
        tbl_element.insert(0, tblPr)
    for el in list(tblPr):
        if el.tag.endswith('}tblStyle') or el.tag.endswith('}tblLook'):
            tblPr.remove(el)
    # Apply the built-in no-border table style, then override with explicit borders
    style_el = OxmlElement('w:tblStyle')
    style_el.set(qn('w:val'), 'TableNormal')
    tblPr.append(style_el)
    # No vertical rules; keep horizontal rules
    borders = OxmlElement('w:tblBorders')
    # Horizontal rules only; omit vertical border elements entirely
    for border_name in ('top', 'bottom', 'insideH'):
        border_el = OxmlElement(f'w:{border_name}')
        border_el.set(qn('w:val'), 'single')
        border_el.set(qn('w:sz'), '4')
        border_el.set(qn('w:space'), '0')
        border_el.set(qn('w:color'), '000000')
        borders.append(border_el)
    tblPr.append(borders)
    def _set_cell_borders(cell, top='single', bottom='single'):
        tcPr = cell._tc.get_or_add_tcPr()
        for child in list(tcPr):
            if child.tag.endswith('}tcBorders'):
                tcPr.remove(child)
        tcBorders = OxmlElement('w:tcBorders')
        for edge, val in [('top', top), ('bottom', bottom)]:
            el = OxmlElement(f'w:{edge}')
            el.set(qn('w:val'), val)
            el.set(qn('w:sz'), '4')
            el.set(qn('w:space'), '0')
            el.set(qn('w:color'), '000000')
            tcBorders.append(el)
        tcPr.append(tcBorders)

    # Header — math-aware formatting for Greek/subscript in column names
    hdr_row = tbl.rows[0]
    trPr = hdr_row._tr.get_or_add_trPr()
    trHeight = OxmlElement('w:trHeight')
    trHeight.set(qn('w:val'), '500')
    trHeight.set(qn('w:hRule'), 'atLeast')
    trPr.append(trHeight)
    for j, col in enumerate(df.columns):
        c = hdr_row.cells[j]
        _set_cell_borders(c)
        c.text = ""
        p = c.paragraphs[0]
        _emit_math_segment(p, str(col), size=font_size, force_bold=True)
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, v in enumerate(row):
            c = tbl.rows[i].cells[j]
            _set_cell_borders(c)
            c.text = ""
            p = c.paragraphs[0]
            _emit_math_segment(p, str(v), size=font_size)
    if col_widths:
        for row in tbl.rows:
            for idx, w in enumerate(col_widths):
                row.cells[idx].width = w


def add_table_block(doc, title, df, caption, widths=None):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(6)
    r = p.add_run(title)
    set_font(r, size=10, bold=True)
    r2 = p.add_run(f"  {caption}")
    set_font(r2, size=10)
    add_dataframe_as_table(doc, df, col_widths=widths)


# ---------------------------------------------------------------------------
# Build docx manuscript from markdown, replacing figure/table placeholders

INSERT_FIG_RE = re.compile(r"^\*\*\[(Insert |)[Ff]ig(ure|\.)? ?(\d+)[^\]]*\]\*\*$")
INSERT_FIG_JA_RE = re.compile(r"^\*\*［図\s*(\d+)[^］]*］\*\*$")
INSERT_TAB_RE = re.compile(r"^\*\*\[(Insert |)[Tt]able ?(\d+[a-z]?)[^\]]*\]\*\*$")
INSERT_TAB_JA_RE = re.compile(r"^\*\*［表\s*(\d+[a-z]?)[^］]*］\*\*$")


def build_manuscript(lang: str, md_name: str = None, out_name: str = None):
    md_path = os.path.join(MS, md_name or f"manuscript_{lang}.md")
    with open(md_path, encoding="utf-8") as fh:
        lines = fh.readlines()

    # Load tables
    t1 = pd.read_csv(os.path.join(TAB, "table1_model_metrics.csv"))
    t2 = pd.read_csv(os.path.join(TAB, "table2_correspondence.csv"))
    t3_path = os.path.join(TAB, "table3_rpim.csv")
    t3 = pd.read_csv(t3_path) if os.path.exists(t3_path) else None
    t4_path = os.path.join(TAB, "table4_extended_oos.csv")
    t4 = pd.read_csv(t4_path) if os.path.exists(t4_path) else None
    t5_path = os.path.join(TAB, "table5_k_level.csv")
    t5 = pd.read_csv(t5_path) if os.path.exists(t5_path) else None
    t6_path = os.path.join(TAB, "table6_tempo_artifact.csv")
    t6 = pd.read_csv(t6_path) if os.path.exists(t6_path) else None
    t7_path = os.path.join(TAB, "table7_solow_episodes.csv")
    t7 = pd.read_csv(t7_path) if os.path.exists(t7_path) else None
    t7b_path = os.path.join(TAB, "table7b_narrative_loo.csv")
    t7b = pd.read_csv(t7b_path) if os.path.exists(t7b_path) else None
    t8_path = os.path.join(TAB, "table8_counterfactual_narrative.csv")
    t8 = pd.read_csv(t8_path) if os.path.exists(t8_path) else None
    t9_path = os.path.join(TAB, "table9_future_scenarios.csv")
    t9 = pd.read_csv(t9_path) if os.path.exists(t9_path) else None
    t10_path = os.path.join(TAB, "table10_cluster_analysis.csv")
    t10 = pd.read_csv(t10_path) if os.path.exists(t10_path) else None
    t10b_path = os.path.join(TAB, "table10b_cluster_metrics.csv")
    t10b = pd.read_csv(t10b_path) if os.path.exists(t10b_path) else None
    t10c_path = os.path.join(TAB, "table10c_cluster_diagnostics.csv")
    t10c = pd.read_csv(t10c_path) if os.path.exists(t10c_path) else None
    t11_path = os.path.join(TAB, "table11_delta_timevarying.csv")
    t11 = pd.read_csv(t11_path) if os.path.exists(t11_path) else None
    t12_path = os.path.join(TAB, "table12_monte_carlo.csv")
    t12 = pd.read_csv(t12_path) if os.path.exists(t12_path) else None
    t13_path = os.path.join(TAB, "table13_asset_lag_robustness.csv")
    t13 = pd.read_csv(t13_path) if os.path.exists(t13_path) else None
    t14_path = os.path.join(TAB, "table14_gamma_price_summary.csv")
    t14 = pd.read_csv(t14_path) if os.path.exists(t14_path) else None

    # Figure caption lookup by index
    fig_cap = {}
    for key, en_prefix, ja_prefix, en_cap, ja_cap, pattern in FIG_LIST:
        idx = int(key.replace("fig", ""))
        fig_cap[idx] = {
            "en": (en_prefix, en_cap, pattern.format(lang="en")),
            "ja": (ja_prefix, ja_cap, pattern.format(lang="ja")),
        }

    t1_cap_en = "Population-capital tempo correspondence."
    t1_cap_ja = "人口・資本テンポ対応関係。"
    t2_cap_en = "M0–M4: in-sample and out-of-sample performance across 39 countries. "\
                "Medians across countries; IQR in brackets."
    t2_cap_ja = "M0–M4: 39カ国の標本内・標本外パフォーマンス（国間中央値、IQRを括弧内）。"
    t3_cap_en = "Relational PIM diagnostics: rho_2 summary under M0, M1, M2, M4."
    t3_cap_ja = "関係型PIM診断: M0, M1, M2, M4におけるρ̂₂の要約。"
    t4_cap_en = "Extended OOS metrics: direction accuracy and CWON trajectory RMSE."
    t4_cap_ja = "拡張標本外指標: 方向精度およびCWON軌跡RMSE。"
    t5_cap_en = "K-level measurement consequences: K gap, TFP shift (= TFP_M0 - TFP_obs), and implied labour-share shift (2010–2019 country means)."
    t5_cap_ja = "K水準の計測帰結: K乖離、TFPシフト（= TFP_M0 - TFP_obs）、労働分配率シフト（2010–2019年国別平均）。"
    t6_cap_en = "Tempo-artifact share of TFP-growth variance: percentage reduction in Var(d log TFP) from M0 to M2 (tempo) and M0 to M4 (joint)."
    t6_cap_ja = "テンポ・アーティファクトのTFP成長率分散シェア: M0→M2（テンポ）およびM0→M4（統合）。"
    t7_cap_en = "Selected historical episodes: mean TFP growth under M0, M2, and M4 (percentage points per year)."
    t7b_cap_en = "Leave-one-out robustness of selected historical episodes: episode removed, and mean TFP growth under M0, M2, and M4."
    t7_cap_ja = "主要な歴史的エピソード: M0, M2, M4下の平均TFP成長率（年率パーセントポイント）。"
    t7b_cap_ja = "主要な歴史的エピソードの抜き一本頑健性：除外したエピソードとM0, M2, M4下の平均TFP成長率。"
    t8_cap_en = "Counterfactual wealth adjustments when the intangible share beta enters official statistics, with 1.5×–3× R&D-based intangible proxies (top countries)."
    t8_cap_ja = "無形資本シェアβを公式統計に含めた場合の反事実国富調整：R&Dベース無形資本プロキシを1.5倍・2倍・3倍とした感度（上位国）。"
    t9_cap_en = "Future GDP-level scenarios, 2020–2040: 2040 and 2030 GDP indices (2019=100), including M4 baseline, level-only TFP, AI surge, and fiscal-stimulus variants."
    t9_cap_ja = "将来GDP水準シナリオ（2020–2040年）：2040年・2030年のGDP指数（2019年=100）。M4ベースライン、TFPレベル固定、AI急増、財政刺激バリエーションを含む。"
    t10_cap_en = "Cross-country clusters by joint-identified tempo lag, intangible share, R&D intensity, and asset composition."
    t10b_cap_en = "Cluster validation metrics: silhouette score and bootstrap stability across candidate k (2–6)."
    t10c_cap_en = "Cluster diagnostics: per-country cluster assignment and silhouette contribution."
    t10_cap_ja = "統合同定されたテンポラグ、無形資本シェア、R&D強度、資産構成による国クラスター。"
    t10b_cap_ja = "クラスター妥当性指標：候補k（2–6）別のシルエット値とブートストラップ安定性。"
    t10c_cap_ja = "クラスター診断：国別クラスター割当てとシルエット寄与。"
    t11_cap_en = "Time-varying depreciation robustness: median estimated lag across countries under alternative delta(t) assumptions."
    t11_cap_ja = "時変減価償却率の頑健性: 代替的δ(t)仮定下の国間中央値推定ラグ。"
    t12_cap_en = "Monte Carlo identification sharpness for four calibrated economies (United States, Republic of Korea, France, Colombia): RMSE of joint-identified mu and beta by sample length and output noise."
    t12_cap_ja = "4カ国（米国、韓国、フランス、コロンビア）のモンテカルロ同定鋭度: サンプル長と出力ショック別の統合同定μ・βのRMSE。"
    t13_cap_en = "M_obs asset-lag robustness: OOS MAPE under scaled literature-based lags."
    t13_cap_ja = "M_obsの資産別懐胎ラグ頑健性：文献ベースのラグをスケーリングした場合の標本外MAPE。"
    t14_cap_en = "Cross-country γ_price summary: zero-crossing price-revaluation rate that would close the PIM-CWON gap, with in-range flag."
    t14_cap_ja = "国別γ_price要約：PIM-CWONギャップを閉じる価格再評価率の零点と範囲内フラグ。"

    # Replace dynamic placeholders from reproducible outputs
    placeholders = {}
    oos_summary_path = os.path.join(DATA, "oos_summary.json")
    if os.path.exists(oos_summary_path):
        with open(oos_summary_path, encoding="utf-8") as fh:
            oos_summary = json.load(fh)
        placeholders["__WILCOXON_P__"] = f"{oos_summary.get('m2_vs_m0_wilcoxon_p', float('nan')):.2g}"
        placeholders["__WILCOXON_N__"] = str(oos_summary.get("m2_vs_m0_n_pairs", 39))
    cond_oos_path = os.path.join(DATA, "conditional_oos.json")
    if os.path.exists(cond_oos_path):
        with open(cond_oos_path, encoding="utf-8") as fh:
            cond_oos = json.load(fh)
        n_boundary = cond_oos.get("n_boundary", 25)
        n_interior = cond_oos.get("n_interior", 14)
        n_countries = n_boundary + n_interior
        placeholders["__N_BOUNDARY__"] = str(n_boundary)
        placeholders["__N_INTERIOR__"] = str(n_interior)
        placeholders["__N_COUNTRIES__"] = str(n_countries)
        fair_path = os.path.join(DATA, "fair_eval.csv")
        if os.path.exists(fair_path) and pd is not None:
            fair = pd.read_csv(fair_path)
            if "country" in fair.columns and "mu_M1" in fair.columns:
                mu = fair.set_index("country")["mu_M1"]
                placeholders["__N_LOWER_BOUND__"] = str(int((mu <= 0.02).sum()))
                placeholders["__N_UPPER_BOUND__"] = str(int((mu >= 5.9).sum()))
    if placeholders:
        lines = [
            functools.reduce(lambda s, kv: s.replace(kv[0], kv[1]), placeholders.items(), line)
            for line in lines
        ]

    doc = Document()
    # Remove identifying author metadata for double-anonymized submission
    doc.core_properties.author = ""
    # Use sensible page margins
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # For double-anonymized submission, the title page (author info,
    # declarations) is built as a separate file; do not include it in the
    # anonymized manuscript body.

    section_num = 0
    subsection_num = 0
    in_references = False
    in_title_page = False

    # Sequential display counters for supplementary figures/tables.
    # The placeholder number is treated as the content key; the displayed
    # number increments in order of appearance.
    supp_fig_count = 0
    supp_table_count = 0

    i = 0
    while i < len(lines):
        line = lines[i].rstrip("\n")
        stripped = line.strip()
        if not stripped:
            i += 1
            continue

        # Title-page block is uploaded separately; skip it entirely.
        # Check end markers first so the Japanese start marker (a prefix of the
        # end marker) does not swallow the end marker.
        if stripped.startswith("[END TITLE PAGE") or stripped.startswith("[タイトルページ終了"):
            in_title_page = False
            i += 1
            continue
        if stripped.startswith("[TITLE PAGE") or stripped.startswith("[タイトルページ"):
            in_title_page = True
            i += 1
            continue
        if in_title_page:
            i += 1
            continue

        # Skip manuscript-body marker line
        if stripped.startswith("[MANUSCRIPT") or stripped.startswith("[原稿"):
            i += 1
            continue

        # Headings
        if stripped.startswith("# "):
            # Paper title (level 1) — keep title case
            add_heading(doc, stripped[2:], 1, lang)
        elif stripped.startswith("## "):
            # Section heading (level 2) — Roman numeral + ALL CAPS
            section_text = stripped[3:]
            m = re.match(r"^(\d+)\s+(.*)$", section_text)
            if m:
                section_num = int(m.group(1))
                section_text = m.group(2)
            else:
                section_num = 0
            in_references = section_text in ("References", "参考文献")
            add_heading(doc, section_text, 2, lang, section_num=section_num)
            subsection_num = 0
        elif stripped.startswith("### "):
            # Subsection heading (level 3) — Roman prefix
            subsection_text = stripped[4:]
            m = re.match(r"^(\d+)\.(\d+)\s+(.*)$", subsection_text)
            if m:
                section_num = int(m.group(1))
                subsection_num = int(m.group(2))
                subsection_text = m.group(3)
            add_heading(doc, subsection_text, 3, lang,
                        section_num=section_num, subsection_num=subsection_num)
        elif stripped.startswith("---"):
            pass  # ignore horizontal rules
        else:
            # Figure inline placeholder (both lang patterns)
            m_fig = INSERT_FIG_RE.match(stripped) or INSERT_FIG_JA_RE.match(stripped)
            m_tab = INSERT_TAB_RE.match(stripped) or INSERT_TAB_JA_RE.match(stripped)
            is_main = md_name is None or ("manuscript" in md_name and "supporting" not in md_name)

            if m_fig:
                try:
                    idx = int(m_fig.group(m_fig.lastindex))
                except Exception:
                    idx = None
                if idx is not None:
                    if is_main and lang == "en":
                        main_fig_display_to_key = {1: 2, 2: 3, 3: 10, 4: 11, 5: 12}
                        key_idx = main_fig_display_to_key.get(idx, idx)
                        display_idx = idx
                        base_prefix = "Fig."
                    elif is_main and lang == "ja":
                        key_idx = idx
                        display_idx = idx
                        base_prefix = "図"
                    else:
                        key_idx = idx
                        supp_fig_count += 1
                        display_idx = supp_fig_count
                        base_prefix = "Supplementary Fig."
                    if key_idx in fig_cap:
                        _prefix, cap, fname = fig_cap[key_idx][lang]
                        prefix = f"{base_prefix} {display_idx}"
                        FIGURE_DISPLAY_MAP[key_idx] = (base_prefix, display_idx)
                        png = os.path.join(FIG, fname)
                        if os.path.exists(png):
                            add_figure(doc, png, prefix, cap)
            elif m_tab:
                idx_str = m_tab.group(m_tab.lastindex)
                try:
                    idx_key = int(idx_str) if idx_str.isdigit() else idx_str
                except Exception:
                    idx_key = None
                if idx_key is None:
                    pass
                elif is_main and lang == "en":
                    if idx_key == 1 and t1 is not None:
                        add_table_block(
                            doc,
                            "Table 1.",
                            t1,
                            t2_cap_en,
                        )
                    elif idx_key == 2 and t5 is not None:
                        add_table_block(
                            doc,
                            "Table 2.",
                            t5,
                            t5_cap_en,
                        )
                elif is_main and lang == "ja":
                    ja_table_map = {
                        1: (t2, t1_cap_ja),
                        2: (t1, t2_cap_ja),
                        3: (t3, t3_cap_ja),
                        4: (t4, t4_cap_ja),
                        5: (t5, t5_cap_ja),
                        6: (t6, t6_cap_ja),
                        7: (t7, t7_cap_ja),
                        8: (t8, t8_cap_ja),
                        9: (t9, t9_cap_ja),
                        10: (t10, t10_cap_ja),
                        11: (t11, t11_cap_ja),
                        12: (t12, t12_cap_ja),
                        13: (t13, t13_cap_ja),
                    }
                    if idx_key in ja_table_map and ja_table_map[idx_key][0] is not None:
                        df, cap = ja_table_map[idx_key]
                        prefix = f"表 {idx_key}."
                        TABLE_DISPLAY_MAP[idx_key] = ("表", idx_key)
                        add_table_block(
                            doc,
                            prefix,
                            df,
                            cap,
                        )
                else:
                    supp_table_map = {
                        1: (t2, t1_cap_en if lang == "en" else t1_cap_ja),
                        3: (t3, t3_cap_en if lang == "en" else t3_cap_ja),
                        4: (t4, t4_cap_en if lang == "en" else t4_cap_ja),
                        5: (t5, t5_cap_en if lang == "en" else t5_cap_ja),
                        6: (t6, t6_cap_en if lang == "en" else t6_cap_ja),
                        7: (t7, t7_cap_en if lang == "en" else t7_cap_ja),
                        "7b": (t7b, t7b_cap_en if lang == "en" else t7b_cap_ja),
                        8: (t8, t8_cap_en if lang == "en" else t8_cap_ja),
                        9: (t9, t9_cap_en if lang == "en" else t9_cap_ja),
                        10: (t10, t10_cap_en if lang == "en" else t10_cap_ja),
                        "10b": (t10b, t10b_cap_en if lang == "en" else t10b_cap_ja),
                        "10c": (t10c, t10c_cap_en if lang == "en" else t10c_cap_ja),
                        11: (t11, t11_cap_en if lang == "en" else t11_cap_ja),
                        12: (t12, t12_cap_en if lang == "en" else t12_cap_ja),
                        13: (t13, t13_cap_en if lang == "en" else t13_cap_ja),
                        14: (t14, t14_cap_en if lang == "en" else t14_cap_ja),
                    }
                    if idx_str in supp_table_map and supp_table_map[idx_str][0] is not None:
                        df, cap = supp_table_map[idx_str]
                        supp_table_count += 1
                        prefix = f"Supplementary Table {supp_table_count}."
                        TABLE_DISPLAY_MAP[idx_str] = ("Supplementary Table", supp_table_count)
                        add_table_block(
                            doc,
                            prefix,
                            df,
                            cap,
                        )
            else:
                # Displayed equation: 4-space indent with label (M0), (1), etc.
                m_eq = EQUATION_RE.match(line)
                if m_eq:
                    add_equation_block(doc, m_eq.group(1), m_eq.group(2))
                    i += 1
                    continue

                text = stripped
                # section separator
                if text in ("Tables", "表", "References", "参考文献"):
                    in_references = text in ("References", "参考文献")
                    add_heading(doc, text, 2, lang)
                # Bullet list items: * text... or - text...
                elif text.startswith("* "):
                    add_rich_para(doc, "•  " + text[2:], lang, bullet=True)
                elif text.startswith("- "):
                    add_rich_para(doc, "•  " + text[2:], lang, bullet=True)
                else:
                    if in_references:
                        text = format_elsevier_reference(text)
                    if lang == "en":
                        text = preserve_en_crossrefs(text)
                    add_rich_para(doc, text, lang)
        i += 1

    # Add centred page numbers to every section
    for section in doc.sections:
        add_page_number_footer(section)

    out = os.path.join(MS, f"{out_name or f'manuscript_{lang}'}.docx")
    doc.save(out)
    print("wrote", out)


# ---------------------------------------------------------------------------
# Separate table .docx files

def build_standalone_tables():
    # Remove stale standalone table docx files before regenerating
    for stale in glob.glob(os.path.join(MS, "table*.docx")):
        os.remove(stale)
    t1 = pd.read_csv(os.path.join(TAB, "table1_model_metrics.csv"))
    t2 = pd.read_csv(os.path.join(TAB, "table2_correspondence.csv"))
    t3_path = os.path.join(TAB, "table3_rpim.csv")
    t3 = pd.read_csv(t3_path) if os.path.exists(t3_path) else None
    table_list = [
        ("table1_model_metrics.docx", t1,
         "Table 1. M0–M4: in-sample and out-of-sample performance across 39 countries. "
         "Medians across countries; IQR in brackets.",
         None),
        ("table2_correspondence.docx", t2,
         "Table 2. Population-capital tempo correspondence.",
         [Inches(1.3), Inches(2.5), Inches(2.5)]),
    ]
    if t3 is not None:
        table_list.append(
            ("table3_rpim.docx", t3,
             "Table 3. Relational PIM diagnostics: rho_2 summary under M0, M1, M2, M4.",
             None))
    t4_path = os.path.join(TAB, "table4_extended_oos.csv")
    t4 = pd.read_csv(t4_path) if os.path.exists(t4_path) else None
    if t4 is not None:
        table_list.append(
            ("table4_extended_oos.docx", t4,
             "Table 4. Extended OOS metrics: direction accuracy and CWON trajectory RMSE.",
             None))
    t5_path = os.path.join(TAB, "table5_k_level.csv")
    t5 = pd.read_csv(t5_path) if os.path.exists(t5_path) else None
    if t5 is not None:
        table_list.append(
            ("table5_k_level.docx", t5,
             "Table 5. K-level measurement consequences: K gap, TFP shift (= TFP_M0 - TFP_obs), and implied labour-share shift (2010–2019 country means).", 
             None))
    t6_path = os.path.join(TAB, "table6_tempo_artifact.csv")
    t6 = pd.read_csv(t6_path) if os.path.exists(t6_path) else None
    if t6 is not None:
        table_list.append(
            ("table6_tempo_artifact.docx", t6,
             "Table 6. Tempo-artifact share of TFP-growth variance.",
             None))
    t7_path = os.path.join(TAB, "table7_solow_episodes.csv")
    t7 = pd.read_csv(t7_path) if os.path.exists(t7_path) else None
    if t7 is not None:
        table_list.append(
            ("table7_solow_episodes.docx", t7,
             "Table 7. Selected historical episodes: mean TFP growth under M0, M2, and M4.",
             None))
    t7b_path = os.path.join(TAB, "table7b_narrative_loo.csv")
    t7b = pd.read_csv(t7b_path) if os.path.exists(t7b_path) else None
    if t7b is not None:
        table_list.append(
            ("table7b_narrative_loo.docx", t7b,
             "Table 7b. Leave-one-out robustness of selected historical episodes.",
             None))
    t8_path = os.path.join(TAB, "table8_counterfactual_narrative.csv")
    t8 = pd.read_csv(t8_path) if os.path.exists(t8_path) else None
    if t8 is not None:
        table_list.append(
            ("table8_counterfactual_narrative.docx", t8,
             "Table 8. Counterfactual wealth adjustments when the intangible share beta enters official statistics, with 1.5×–3× R&D-based intangible proxies.",
             None))
    t9_path = os.path.join(TAB, "table9_future_scenarios.csv")
    t9 = pd.read_csv(t9_path) if os.path.exists(t9_path) else None
    if t9 is not None:
        table_list.append(
            ("table9_future_scenarios.docx", t9,
             "Table 9. Future GDP-level scenarios, 2020–2040: 2040 and 2030 GDP indices (2019=100), including M4 baseline, level-only TFP, AI surge, and fiscal-stimulus variants.",
             None))
    t10_path = os.path.join(TAB, "table10_cluster_analysis.csv")
    t10 = pd.read_csv(t10_path) if os.path.exists(t10_path) else None
    if t10 is not None:
        table_list.append(
            ("table10_cluster_analysis.docx", t10,
             "Table 10. Cross-country clusters by joint-identified tempo lag, intangible share, R&D intensity, and asset composition.",
             None))
    t10b_path = os.path.join(TAB, "table10b_cluster_metrics.csv")
    t10b = pd.read_csv(t10b_path) if os.path.exists(t10b_path) else None
    if t10b is not None:
        table_list.append(
            ("table10b_cluster_metrics.docx", t10b,
             "Table 10b. Cluster validation metrics (silhouette and bootstrap stability).",
             None))
    t10c_path = os.path.join(TAB, "table10c_cluster_diagnostics.csv")
    t10c = pd.read_csv(t10c_path) if os.path.exists(t10c_path) else None
    if t10c is not None:
        table_list.append(
            ("table10c_cluster_diagnostics.docx", t10c,
             "Table 10c. Cluster diagnostics: per-country assignment and silhouette contribution.",
             None))
    t11_path = os.path.join(TAB, "table11_delta_timevarying.csv")
    t11 = pd.read_csv(t11_path) if os.path.exists(t11_path) else None
    if t11 is not None:
        table_list.append(
            ("table11_delta_timevarying.docx", t11,
             "Table 11. Time-varying depreciation robustness: median estimated lag across countries.",
             None))
    t12_path = os.path.join(TAB, "table12_monte_carlo.csv")
    t12 = pd.read_csv(t12_path) if os.path.exists(t12_path) else None
    if t12 is not None:
        table_list.append(
            ("table12_monte_carlo.docx", t12,
             "Table 12. Monte Carlo identification sharpness for four calibrated economies (United States, Republic of Korea, France, Colombia).",
             None))
    t13_path = os.path.join(TAB, "table13_asset_lag_robustness.csv")
    t13 = pd.read_csv(t13_path) if os.path.exists(t13_path) else None
    if t13 is not None:
        table_list.append(
            ("table13_asset_lag_robustness.docx", t13,
             "Table 13. M_obs asset-lag robustness: OOS MAPE under scaled literature-based lags.",
             None))
    t14_path = os.path.join(TAB, "table14_gamma_price_summary.csv")
    t14 = pd.read_csv(t14_path) if os.path.exists(t14_path) else None
    if t14 is not None:
        table_list.append(
            ("table14_gamma_price_summary.docx", t14,
             "Table 14. Cross-country γ_price summary: zero-crossing price-revaluation rate that closes the PIM-CWON gap.",
             None))
    for name, df, cap, widths in table_list:
        d = Document()
        add_heading(d, cap, 2, "en")
        add_dataframe_as_table(d, df, col_widths=widths)
        out = os.path.join(MS, name)
        d.save(out)
        print("wrote", out)


# ---------------------------------------------------------------------------
# Editable pptx (English figures only, one per slide)

def build_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    for key, en_prefix, _ja, en_cap, _jc, pattern in FIG_LIST:
        png = os.path.join(FIG, pattern.format(lang="en"))
        if not os.path.exists(png):
            print("skip missing", png)
            continue
        idx = int(key.replace("fig", ""))
        mapping = FIGURE_DISPLAY_MAP.get(idx)
        if mapping is None:
            parts = en_prefix.split()
            prefix = parts[0]
            display_num = int(parts[-1]) if parts[-1].isdigit() else idx
        else:
            prefix, display_num = mapping
        label = f"{prefix} {display_num}"
        slide = prs.slides.add_slide(blank)
        # Title
        tb = slide.shapes.add_textbox(
            PptInches(0.4), PptInches(0.2),
            PptInches(12.5), PptInches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{label}. {en_cap}"
        r.font.size = PptPt(20)
        r.font.bold = True
        # Image: fit to 12 x 5.5 area, centered
        slide.shapes.add_picture(
            png, PptInches(1.5), PptInches(1.1),
            width=PptInches(10.3), height=PptInches(5.6))
        # Caption bar at bottom
        cb = slide.shapes.add_textbox(
            PptInches(0.4), PptInches(6.85),
            PptInches(12.5), PptInches(0.5))
        cf = cb.text_frame
        cf.word_wrap = True
        pp = cf.paragraphs[0]
        rr = pp.add_run()
        rr.text = en_cap
        rr.font.size = PptPt(14)
        rr.font.italic = True
        rr.font.color.rgb = PptRGB(0x44, 0x44, 0x44)

    out = os.path.join(MS, "figures_en.pptx")
    prs.save(out)
    print("wrote", out)


# ---------------------------------------------------------------------------
# Separate figure files for submission (Wiley/Economica preferred)

def build_separate_figures():
    """Copy English figure PNGs to manuscript/figures/ for submission.

    Main manuscript figures are named Figure_1.png ... Figure_5.png;
    supplementary figures are named Supplementary_Figure_1.png ...
    according to their sequential display numbers.
    """
    sep_dir = os.path.join(MS, "figures")
    os.makedirs(sep_dir, exist_ok=True)
    # Remove stale generated figure names before rewriting them
    for stale in glob.glob(os.path.join(sep_dir, "Figure_*.png")) + glob.glob(os.path.join(sep_dir, "Supplementary_Figure_*.png")):
        os.remove(stale)
    for key, en_prefix, _ja, _en_cap, _ja_cap, pattern in FIG_LIST:
        idx = int(key.replace("fig", ""))
        src = os.path.join(FIG, pattern.format(lang="en"))
        if not os.path.exists(src):
            print("skip missing separate figure", src)
            continue
        mapping = FIGURE_DISPLAY_MAP.get(idx)
        if mapping is None:
            parts = en_prefix.split()
            prefix = parts[0]
            display_num = int(parts[-1]) if parts[-1].isdigit() else idx
        else:
            prefix, display_num = mapping
        if "Supplementary" in prefix:
            dst = os.path.join(sep_dir, f"Supplementary_Figure_{display_num}.png")
        else:
            dst = os.path.join(sep_dir, f"Figure_{display_num}.png")
        shutil.copy2(src, dst)
        print(f"wrote {dst}")


# ---------------------------------------------------------------------------
# PDF conversion via LibreOffice (fonts embedded by default)

def convert_docx_to_pdf(docx_path: str):
    """Convert a .docx file to PDF using LibreOffice.

    LibreOffice embeds fonts by default, satisfying the Editorial Express
    requirement that all fonts be embedded in the PDF.
    """
    import subprocess
    executable = shutil.which("libreoffice")
    if executable is None:
        print(f"SKIP: LibreOffice is unavailable; not rebuilding {docx_path[:-5]}.pdf")
        return
    out_dir = os.path.dirname(docx_path)
    cmd = [
        executable, "--headless", "--convert-to", "pdf",
        "--outdir", out_dir, docx_path,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        print(f"WARNING: PDF conversion failed for {docx_path}")
        print(result.stderr)
    else:
        pdf_path = docx_path.replace(".docx", ".pdf")
        print(f"wrote {pdf_path}")


def build_title_page():
    """Build separate title page .docx with author info (for double-blind submission)."""
    md_path = os.path.join(MS, "manuscript_en.md")
    if not os.path.exists(md_path):
        return
    with open(md_path, encoding="utf-8") as fh:
        lines = fh.readlines()

    # Extract title-page content between markers
    title_lines = []
    in_tp = False
    for raw_line in lines:
        line = raw_line.rstrip("\n").strip()
        if line.startswith("[TITLE PAGE"):
            in_tp = True
            continue
        if line.startswith("[END TITLE PAGE"):
            break
        if in_tp:
            title_lines.append(raw_line.rstrip("\n"))

    if not title_lines:
        print("no title page content found")
        return

    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    for raw_line in title_lines:
        stripped = raw_line.strip()
        if not stripped:
            spacer = doc.add_paragraph()
            spacer.paragraph_format.space_after = Pt(0)
            spacer.paragraph_format.space_before = Pt(0)
            spacer.paragraph_format.line_spacing = 1.5
            continue
        if stripped.startswith("# "):
            add_heading(doc, stripped[2:], 1, "en")
        else:
            add_rich_para(doc, stripped, "en")

    out = os.path.join(MS, "title_page_en.docx")
    doc.save(out)
    print(f"wrote {out}")
    convert_docx_to_pdf(out)


def build_cover_letter():
    """Build cover letter docx from markdown source."""
    CL = os.path.join(ROOT, "cover_letter")
    md_path = os.path.join(CL, "cover_letter_en.md")
    if not os.path.exists(md_path):
        print(f"cover letter not found: {md_path}")
        return
    with open(md_path, encoding="utf-8") as fh:
        lines = fh.readlines()

    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    for raw_line in lines:
        line = raw_line.rstrip("\n")
        stripped = line.strip()

        if not stripped:
            # Blank line → small spacer paragraph
            spacer = doc.add_paragraph()
            spacer.paragraph_format.space_after = Pt(0)
            spacer.paragraph_format.space_before = Pt(0)
            spacer.paragraph_format.line_spacing = 1.5
            continue

        # Bullet list items
        if stripped.startswith("- "):
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(2)
            p.paragraph_format.line_spacing = 1.5
            p.paragraph_format.left_indent = Pt(36)
            p.paragraph_format.first_line_indent = Pt(-18)
            text = "•  " + stripped[2:]
            text = re.sub(r'\\([*_\\`])', r'\1', text)
            add_math_runs(p, text, size=10, base_italic=False)
            continue

        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(2)
        p.paragraph_format.line_spacing = 1.5

        # Handle bold + italic markdown
        text = stripped
        text = re.sub(r'\\([*_\\`])', r'\1', text)

        # Split on **bold** markers
        parts = re.split(r'(\*\*[^*]+\*\*)', text)
        for part in parts:
            if not part:
                continue
            if part.startswith('**') and part.endswith('**') and len(part) > 4:
                inner = part[2:-2]
                # Bold — further split on *italic*
                sub_parts = re.split(r'(\*[^*]+\*)', inner)
                for sp in sub_parts:
                    if not sp:
                        continue
                    if sp.startswith('*') and sp.endswith('*') and len(sp) > 2:
                        _emit_math_segment(p, sp[1:-1], size=10,
                                           force_italic=True, force_bold=True)
                    else:
                        _emit_math_segment(p, sp, size=10,
                                           force_italic=False, force_bold=True)
            else:
                # Normal text — handle *italic* within
                add_math_runs(p, part, size=10, base_italic=False)

    out = os.path.join(CL, "cover_letter_en.docx")
    doc.save(out)
    print(f"wrote {out}")
    convert_docx_to_pdf(out)


def build_highlights():
    """Build highlights docx from markdown source (3-5 bullets, <=85 chars)."""
    md_path = os.path.join(MS, "highlights_en.md")
    if not os.path.exists(md_path):
        print(f"highlights not found: {md_path}")
        return
    with open(md_path, encoding="utf-8") as fh:
        lines = fh.readlines()

    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(2)
    p.paragraph_format.line_spacing = 1.5
    run = p.add_run("Highlights")
    set_font(run, name="Times New Roman", size=12, bold=True)

    for raw_line in lines:
        line = raw_line.rstrip("\n").strip()
        if not line:
            continue
        if line.startswith(("- ", "* ")):
            line = line[2:]
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(2)
        p.paragraph_format.line_spacing = 1.5
        p.paragraph_format.left_indent = Pt(36)
        p.paragraph_format.first_line_indent = Pt(-18)
        run = p.add_run("\u2022  ")
        set_font(run, name="Times New Roman", size=10, bold=False)
        add_math_runs(p, line, size=10, base_italic=False)

    out = os.path.join(MS, "highlights_en.docx")
    doc.save(out)
    print(f"wrote {out}")
    convert_docx_to_pdf(out)


def main():
    # English version determines figure numbering for the submission PPTX and
    # separate figure files, so build it (and its SI) before those outputs.
    build_manuscript("en")
    if os.path.exists(os.path.join(MS, "supporting_information_en.md")):
        build_manuscript("en", md_name="supporting_information_en.md", out_name="supporting_information_en")
    build_standalone_tables()
    build_pptx()
    build_separate_figures()
    build_title_page()
    build_highlights()
    build_cover_letter()
    build_manuscript("ja")
    # Convert manuscripts to PDF
    for name in ("manuscript_en", "manuscript_ja", "supporting_information_en", "highlights_en"):
        docx = os.path.join(MS, f"{name}.docx")
        if os.path.exists(docx):
            convert_docx_to_pdf(docx)


if __name__ == "__main__":
    main()
